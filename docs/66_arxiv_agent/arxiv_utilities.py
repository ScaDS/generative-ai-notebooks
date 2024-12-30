import requests
import xml.etree.ElementTree as ET
from IPython.display import display, Markdown

def search_arxiv(query=None, author=None, year=None, max_results=1):
    # Base URL for the arXiv API
    base_url = "http://export.arxiv.org/api/query"

    # Construct the query string
    query_parts = []
    if query:
        query_parts.append(f"all:{query}")
    if author:
        query_parts.append(f"au:{author}")
    if year:
        query_parts.append(f"submittedDate:[{year}01010000 TO {year}12312359]")
    
    # Combine all parts into a single query
    full_query = " AND ".join(query_parts)

    # Construct the query parameters
    params = {
        "search_query": full_query,
        "start": 0,
        "max_results": max_results * 10,
        "sortBy": "submittedDate",
        "sortOrder": "descending"
    }

    # Send the request
    response = requests.get(base_url, params=params)
    response.raise_for_status()

    #print("Response:", response.text)

    # Parse the XML response
    root = ET.fromstring(response.text)

    #
    papers = []
    for entry in root.findall("{http://www.w3.org/2005/Atom}entry"):
        title = entry.find("{http://www.w3.org/2005/Atom}title").text.strip()
        authors = [author.find("{http://www.w3.org/2005/Atom}name").text for author in entry.findall("{http://www.w3.org/2005/Atom}author")]
        published = entry.find("{http://www.w3.org/2005/Atom}published").text
        summary = entry.find("{http://www.w3.org/2005/Atom}summary").text.strip()
        link = entry.find("{http://www.w3.org/2005/Atom}id").text

        # Check the license by elaborating if the CC-BY 4.0 button is on the arxiv abs page
        if check_if_arxiv_paper_is_cc_by_40(link):
            papers.append({
                "title": title,
                "authors": authors,
                "published": published,
                "summary": summary,
                "link": link,
            })

        if len(papers) >= max_results:
            break

    return papers

def check_if_arxiv_paper_is_cc_by_40(link):
    """Returns if an arxiv paper as specified by its link, is CC-BY 4.0 licensed"""
    import requests
    license_response = requests.get(link)        
    return license_response.status_code == 200 and "https://arxiv.org/icons/licenses/by-4.0.png" in license_response.text

def get_arxiv_metadata(arxiv_id):
    import arxiv
    search = arxiv.Search(id_list=[arxiv_id])
    paper = next(search.results())
    
    metadata = {
        'title': paper.title,
        'authors': [author.name for author in paper.authors],
        'published': paper.published,
        'summary': paper.summary,
        'doi': paper.doi,
        'primary_category': paper.primary_category,
        'categories': paper.categories,
        'pdf_url': paper.pdf_url
    }
    return metadata
    
def convert_to_markdown(papers):
    md_content = "# Search Results\n\n"
    for idx, paper in enumerate(papers, start=1):
        md_content += f"## Paper {idx}\n"
        md_content += f"**Title:** {paper['title']}\n\n"
        md_content += f"**Authors:** {', '.join(paper['authors'])}\n\n"
        md_content += f"**Published:** {paper['published']}\n\n"
        md_content += f"**Summary:** {paper['summary']}\n\n"
        md_content += f"**Link:** [Read More]({paper['link']})\n\n"
    return md_content

def download_pdf(paper_link):
    """Downloads an arxiv paper specified by the arxiv-abs-link."""
    # Check license
    if check_if_arxiv_paper_is_cc_by_40(paper_link):
        pdf_url = paper_link.replace("abs", "pdf")

        pdf_response = requests.get(pdf_url)
        if pdf_response.status_code == 200:
            pdf_filename = paper_link.split("/")[-1] + ".pdf"
            with open(pdf_filename, "wb") as pdf_file:
                pdf_file.write(pdf_response.content)
            display(Markdown(f"PDF downloaded: [{paper_link}](pdf_filename), licensed CC-BY 4.0"))
            return pdf_filename
        else:
            print("Failed to download PDF.")
    else:
        print("The paper is not licensed under CC-BY 4.0.")

def pdf_to_markdown(pdf_path):
    """
    Convert a PDF file to a Markdown string.

    Args:
        pdf_path (str): Path to the input PDF file.

    Returns:
        str: A string containing the Markdown representation of the PDF.
    """
    import os
    from pypdf import PdfReader

    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"The file {pdf_path} does not exist.")

    reader = PdfReader(pdf_path)
    markdown_content = []

    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if not text.strip():
            continue

        markdown_content.append(f"\n# Page {page_num}\n\n")  # Markdown header for the page
        lines = text.split('\n')
        for line in lines:
            line = line.strip()
            if line:
                markdown_content.append(f"{line}\n\n")

    return ''.join(markdown_content)


def make_powerpoint_slides(slides:str, filename="slides.pptx"):
    """
    Create a PowerPoint slide deck from a list of slides in JSON format.

    Examples for slides:
    Example 1: {"title":"Continents and backing", "content":["Author1 name, Author2 name, ..."]}
    Example 2: {"title":"Continents", "content":["The planet is separated into continents", "Each continent consists of countries, which might be further structured into federal disctricts"]}
    Example 3: {"title":"Backing cake", "content":["When baking cake it is important to pre-heat the oven before putting the cake in.", "Ingredients: \n  * Flour\n  * Sugar\n  * Salt\n  * Eggs"]}
    Example 4: {"title":"Regional food", "content":["Depending on the continent, regional food cultures have developed over the centures.", "These traditions developed because of availability of different ingredients and also depending on other natural resources.", "For example, in antarctical, wine grows relatively badly. That's why there are no traditional wines from Antarctica and wine needs to be imported in case it is an ingredient for a certain food. Hence, no traditional food from Antarctical is made of wine."]}
    Example 5: {"title":"Summary", "content":["In this slide-deck we learned about\n* Continents\n* Backing\n* Regional differences in food culture"]}

    Parameters
    ----------
    slides_description_json : str
        A JSON string representing the slides. Each slide is expected to be a dictionary with
        a 'title' key and a 'content' key, where 'content' is a list of strings.
    filename : str, optional
        The name of the file to save the presentation to. Defaults to 'issue_slides.pptx'.

    """
    import json
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pathlib import Path
    from PIL import Image
    import os
    import re

    # Parse json-encoded slide description
    slides_description_json = re.sub(r'[\x00-\x1f\x7f]', '', slides)

    slides_data = json.loads(slides_description_json)

    file_location = "slide_template.pptx"
    if not os.path.exists(file_location):
        file_location = Path(__file__).parent / "data" / "slide_template.pptx"

    # Create a presentation
    presentation = Presentation(file_location)

    # determine slide size
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Convert EMU to inches (1 inch = 914400 EMUs)
    width_inch = slide_width / 914400
    height_inch = slide_height / 914400

    top = Inches(2)
    bottom = Inches(1)

    # Iterate through slide data to create slides
    for i, slide_data in enumerate(slides_data):
        # Add a slide with title and content layout
        slide_layout = presentation.slide_layouts[0 if i == 0 else 1] # choose first or second layout
        slide = presentation.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.title
        title_box.text = slide_data['title']

        # Calculate width for content columns
        num_columns = len(slide_data['content'])

        # remove all placeholders except the title
        for shape in slide.placeholders:
            if shape != title_box:
                #shape.text = ""
                top = shape.top
                bottom = height_inch - top
                slide.shapes._spTree.remove(shape._element)

        content = "\n".join(slide_data['content'])
        
        content_width = (Inches(width_inch) - top)
        content_height = (Inches(height_inch) - top) - bottom
        left = Inches(1)

        content_box = slide.shapes.add_textbox(left=left, top=top, width=content_width, height=content_height)
        text_frame = content_box.text_frame
        text_frame.text = content
        text_frame.word_wrap = True

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # Set the font size for each run in each paragraph
                run.font.size = Pt(24)  # Set font size to 24 points


    # Save presentation
    presentation.save(filename)

    return f"The Powerpoint Presentation was saved as [{filename}]({filename})"

def prompt_scadsai_llm(message:str, model="meta-llama/Llama-3.3-70B-Instruct"):
    """A prompt helper function that sends a message to ScaDS.AI LLM server at 
    ZIH TU Dresden and returns only the text response.
    """
    import os
    import openai
    
    # convert message in the right format if necessary
    if isinstance(message, str):
        message = [{"role": "user", "content": message}]
    
    # setup connection to the LLM
    client = openai.OpenAI(base_url="https://llm.scads.ai/v1",
                           api_key=os.environ.get('SCADSAI_API_KEY')
    )
    response = client.chat.completions.create(
        model=model,
        messages=message
    )
    
    # extract answer
    return response.choices[0].message.content
